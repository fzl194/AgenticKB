"""PPTX 原生解析适配器 + Normalizer（M3, SRS §C06/§C07, ADR-0003 D-028A；
2026-08-17 整改轮重写）.

理念：全部用工业级成熟库（python-pptx），自研代码只做"库输出 ->
BackendBlock"映射。

整改轮修正（用户指令）：
- **bbox 修复**：``(x0, top, x1, bottom)`` 角点语义（此前误放
  width/height，跨格式不变量 I-1）。
- **文本框按 paragraph 拆分**：bullet 段落（显式 ``buChar/buAutoNum``
  或 ``lvl>0``）-> ``list_item(level=lvl+1)``，普通段落 -> ``paragraph``；
  title placeholder 首段 -> heading。
- **阅读序不等于 shape XML 顺序**：几何带排序（垂直重叠成带、带间按
  top、带内按 left）；无位置 shape 保持 XML 序置于带序之后（不伪造）。
- **notes 保留**：``notes_slide`` 文本成块（``slide_notes`` 标记，可过滤）。
- **group shape 递归**：组成员按组内坐标参与排序。
- **图片资产**：picture -> figure 块 + sha256 + binary asset 引用，
  Normalizer 产 FigureAsset。
- **chart / SmartArt 诊断**：不支持 -> 计数进 warnings（不静默丢失）。
- 表格 cell 独立证据（evidence_index -> cell 级 span）。

fingerprint：``native_pptx@2.0.0#python-pptx-<ver>``。
"""
from __future__ import annotations

import hashlib
import io
from importlib.metadata import version as _pkg_version
from typing import Any

from pptx import Presentation as OpenPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.base import BaseShape

from knowledge_mining.mining.contracts.parse_ir import (
    Container,
    EvidenceSpan,
    FigureAsset,
)
from knowledge_mining.mining.contracts.parser_adapter import (
    BackendBlock,
    BackendParseArtifact,
    ParserAdapterError,
    ParserDescriptor,
    UnsupportedFormat,
)
from knowledge_mining.mining.parse_adapters.native._base import (
    BaseNativeNormalizer,
)

NATIVE_PPTX_PARSER_ID = "native_pptx"
NATIVE_PPTX_VERSION = "2.0.0"
_PYTHON_PPTX_VERSION = _pkg_version("python-pptx")
NATIVE_PPTX_FINGERPRINT = (
    f"{NATIVE_PPTX_PARSER_ID}@{NATIVE_PPTX_VERSION}"
    f"#python-pptx-{_PYTHON_PPTX_VERSION}"
)

NATIVE_PPTX_MIMES = frozenset({
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
})

_TITLE_PLACEHOLDER_TYPES = frozenset({
    PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
})

# 顶部标题区判定参数（PPT 排版惯例，验收 v9）。
_TITLE_ZONE_RATIO = 0.30   # 页面上部 30% 为标题带
_TITLE_MAX_CHARS = 24      # 标题短句上限（超长是正文/说明框）

_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


class NativePptxParser:
    """DocumentParser 实现：python-pptx 包装（SRS §C06）."""

    def __init__(self) -> None:
        self.descriptor = ParserDescriptor(
            parser_id=NATIVE_PPTX_PARSER_ID,
            display_name="Native PPTX Parser (python-pptx)",
            version=NATIVE_PPTX_VERSION,
            supported_mimes=NATIVE_PPTX_MIMES,
            backend_kind="local",
            parser_fingerprint=NATIVE_PPTX_FINGERPRINT,
            capabilities=frozenset({
                "slides", "titles", "shapes", "tables", "pictures", "notes",
            }),
        )

    def supports(self, mime: str) -> bool:
        return self.descriptor.supports(mime)

    def parse(self, data: bytes, *, mime: str) -> BackendParseArtifact:
        if not self.supports(mime):
            raise UnsupportedFormat(
                f"{NATIVE_PPTX_PARSER_ID} cannot parse mime {mime!r}"
            )
        try:
            prs = OpenPresentation(io.BytesIO(data))
        except Exception as exc:  # 第三方异常不得穿越适配层（SRS §C06）
            raise ParserAdapterError(
                f"{NATIVE_PPTX_PARSER_ID}: python-pptx failed to open source: {exc}"
            ) from exc

        blocks: list[BackendBlock] = []
        warnings: list[str] = []
        try:
            slide_h_emu = float(prs.slide_height)
            for slide_index, slide in enumerate(prs.slides):
                blocks.extend(_slide_blocks(
                    slide, slide_index, slide_h_emu, warnings
                ))
        except Exception as exc:  # 中段损坏也归一（§C06，评审 MED）
            raise ParserAdapterError(
                f"{NATIVE_PPTX_PARSER_ID}: failed to walk slides: {exc}"
            ) from exc

        usage: dict[str, Any] = {
            "slide_width_emu": int(prs.slide_width),
            "slide_height_emu": int(prs.slide_height),
        }
        return BackendParseArtifact(
            parser_id=NATIVE_PPTX_PARSER_ID,
            parser_version=NATIVE_PPTX_VERSION,
            mime=mime.lower(),
            blocks=tuple(blocks),
            usage=usage,
            warnings=tuple(warnings),
        )


# ---------------------------------------------------------------------------
# Parser 映射（模块级纯函数，便于单测）
# ---------------------------------------------------------------------------


def _slide_blocks(
    slide: Any,
    slide_index: int,
    slide_h_emu: float,
    warnings: list[str],
) -> list[BackendBlock]:
    """一页 -> 块序列：几何阅读序的 shape 块 + 末尾 notes 块."""
    entries: list[tuple[tuple[float, float, float, float] | None, Any, list[int]]] = []
    _collect_shapes(
        slide.shapes, slide_index, entries, warnings, path=[]
    )
    has_title_placeholder = any(
        _is_title_placeholder(shape) for _, shape, _ in entries
    )
    ordered = _reading_order(entries)

    blocks: list[BackendBlock] = []
    for bbox, shape, path in ordered:
        blocks.extend(_shape_blocks(
            shape, slide_index, path, bbox, slide_h_emu, warnings,
            has_title_placeholder=has_title_placeholder,
        ))

    notes = _notes_blocks(slide, slide_index)
    blocks.extend(notes)
    return blocks


def _collect_shapes(
    shapes: Any,
    slide_index: int,
    entries: list,
    warnings: list[str],
    path: list[int],
) -> None:
    """收集 shape（group 递归展开；chart/SmartArt 诊断）."""
    for j, shape in enumerate(shapes):
        here = [*path, j]
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _collect_shapes(
                shape.shapes, slide_index, entries, warnings, here
            )
            continue
        if _is_chart(shape):
            warnings.append(
                f"slide {slide_index}: chart shape (path={here}) not "
                "supported (diagnosed, not parsed)"
            )
            continue
        if _is_smartart(shape):
            warnings.append(
                f"slide {slide_index}: SmartArt shape (path={here}) not "
                "supported (diagnosed, not parsed)"
            )
            continue
        entries.append((_bbox_of(shape), shape, here))


def _is_chart(shape: BaseShape) -> bool:
    try:
        return bool(shape.has_chart)
    except Exception:  # noqa: BLE001 —— 非 graphicFrame 无 has_chart 语义
        return False


def _is_smartart(shape: BaseShape) -> bool:
    """graphicFrame 无表无图但带 diagram 数据 -> SmartArt."""
    try:
        if shape.shape_type != MSO_SHAPE_TYPE.GRAPHIC_FRAME:
            return False
        if shape.has_table or shape.has_chart:
            return False
        xml = shape._element.xml  # noqa: SLF001
        return "diagram" in xml
    except Exception:  # noqa: BLE001
        return False


def _reading_order(entries: list) -> list:
    """几何阅读序：垂直重叠成带、带间按 top、带内按 left（整改轮）.

    无 bbox 的 shape 保持 XML 顺序置于末尾（不伪造几何证据）。
    """
    with_box = sorted(
        (e for e in entries if e[0] is not None),
        key=lambda e: (e[0][1], e[0][0]),
    )
    without_box = [e for e in entries if e[0] is None]
    bands: list[dict[str, Any]] = []
    for entry in with_box:
        bbox = entry[0]
        if bands and bbox[1] <= bands[-1]["bottom"] + 1:
            bands[-1]["items"].append(entry)
            bands[-1]["bottom"] = max(bands[-1]["bottom"], bbox[3])
        else:
            bands.append({"bottom": bbox[3], "items": [entry]})
    ordered: list = []
    for band in bands:
        ordered.extend(sorted(band["items"], key=lambda e: e[0][0]))
    ordered.extend(without_box)
    return ordered


def _shape_blocks(
    shape: BaseShape,
    slide_index: int,
    path: list[int],
    bbox: tuple[float, float, float, float] | None,
    slide_h_emu: float,
    warnings: list[str],
    has_title_placeholder: bool = False,
) -> list[BackendBlock]:
    """一个 shape -> 块序列（文本框按段落拆分，整改轮）."""
    common: dict[str, Any] = {
        "container_ref": {"container_type": "slide", "index": slide_index},
        "native_ref": _native_ref(slide_index, path),
        "bbox": bbox,
    }
    if _is_table(shape):
        return [_table_block(shape, common)]
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return [_picture_block(shape, common, warnings)]
    if shape.has_text_frame:
        return _text_frame_blocks(
            shape, common, bbox, slide_h_emu, has_title_placeholder
        )
    return []


def _native_ref(slide_index: int, path: list[int]) -> dict[str, Any]:
    ref: dict[str, Any] = {"slide_index": slide_index, "shape_index": path[0]}
    if len(path) > 1:
        ref["group_path"] = path
    return ref


def _is_table(shape: BaseShape) -> bool:
    try:
        return bool(shape.has_table)
    except Exception:  # noqa: BLE001
        return False


def _text_frame_blocks(
    shape: BaseShape,
    common: dict[str, Any],
    bbox: tuple[float, float, float, float] | None,
    slide_h_emu: float,
    has_title_placeholder: bool = False,
) -> list[BackendBlock]:
    """text_frame -> 逐段落块（bullet/lvl -> list_item，整改轮）.

    标题回退（v9 惯例信号）收紧为三条件同时成立：页面**没有** title
    占位符、框内**只有单个非空段落**、位于页顶标题带且短句——避免把
    多段正文框的首段或已有占位符页面的顶区短正文误升为标题。
    """
    paragraphs = list(shape.text_frame.paragraphs)
    nonempty = [
        p for p in paragraphs
        if "".join(run.text or "" for run in p.runs).strip()
    ]
    single_candidate = (
        len(nonempty) == 1 and not has_title_placeholder
        and not _is_title_placeholder(shape)
        and bbox is not None
        and _looks_like_slide_title(
            "".join(run.text or "" for run in nonempty[0].runs).strip(),
            bbox[1], slide_h_emu,
        )
    )
    is_title = _is_title_placeholder(shape)
    blocks: list[BackendBlock] = []
    for k, p in enumerate(paragraphs):
        text = "".join(run.text or "" for run in p.runs).strip()
        if not text:
            continue
        bullet = _paragraph_bullet_level(p)
        if k == 0 and is_title:
            blocks.append(BackendBlock(
                block_type="heading", text=text, level=1, **common
            ))
            continue
        if single_candidate and not blocks:
            blocks.append(BackendBlock(
                block_type="heading", text=text, level=1, **common
            ))
            continue
        if bullet is not None:
            blocks.append(BackendBlock(
                block_type="list_item", text=text, level=bullet, **common
            ))
        else:
            blocks.append(BackendBlock(
                block_type="paragraph", text=text, **common
            ))
    return blocks


def _paragraph_bullet_level(paragraph: Any) -> int | None:
    """a:pPr -> 列表层级：显式 bullet 符号或 lvl>0 => list_item(lvl+1)."""
    pPr = paragraph._p.find(f"{_A_NS}pPr")  # noqa: SLF001
    if pPr is None:
        return None
    has_bullet = any(
        pPr.find(f"{_A_NS}{tag}") is not None
        for tag in ("buChar", "buAutoNum", "buBlip")
    )
    lvl_attr = pPr.get("lvl")
    try:
        lvl = int(lvl_attr) if lvl_attr is not None else 0
    except ValueError:
        lvl = 0
    if has_bullet or lvl > 0:
        return lvl + 1
    return None


def _picture_block(
    shape: BaseShape,
    common: dict[str, Any],
    warnings: list[str],
) -> BackendBlock:
    """picture -> figure 块 + sha256/binary asset 引用（整改轮）."""
    try:
        image = shape.image
        blob = image.blob
        sha = hashlib.sha256(blob).hexdigest()
        structure = {
            "name": shape.name,
            "image_sha256": sha,
            "binary_asset_id": f"bin-{sha[:16]}",
            "mime": image.content_type,
            "size": len(blob),
        }
    except Exception as exc:  # noqa: BLE001 —— 图片读取失败诊断不中断
        warnings.append(f"picture read failed: {exc}")
        structure = {"name": shape.name}
    return BackendBlock(
        block_type="figure", text="", structure=structure, **common
    )


def _notes_blocks(slide: Any, slide_index: int) -> list[BackendBlock]:
    """演讲备注 -> 块（slide_notes 标记，内容保留可过滤，整改轮）."""
    if not slide.has_notes_slide:
        return []
    text = (slide.notes_slide.notes_text_frame.text or "").strip()
    if not text:
        return []
    return [BackendBlock(
        block_type="paragraph",
        text=text,
        container_ref={"container_type": "slide", "index": slide_index},
        native_ref={"slide_index": slide_index, "notes": True},
        structure={"slide_notes": True},
    )]


def _looks_like_slide_title(text: str, top_emu: float, slide_h: float) -> bool:
    """页顶短行是否像 slide 标题（占位符缺失时的通用回退信号）."""
    if slide_h <= 0:
        return False
    compact = text.strip()
    if not compact or "\n" in compact:
        return False
    return (top_emu / slide_h) < _TITLE_ZONE_RATIO and len(compact) <= _TITLE_MAX_CHARS


def _table_block(shape: BaseShape, common: dict[str, Any]) -> BackendBlock:
    """pptx table -> table 块（is_merge_origin/is_spanned 库 API 直读）."""
    table = shape.table
    rows = len(table.rows)
    cols = len(table.columns)
    cells: list[dict[str, Any]] = []
    evidence = 0
    for r, row in enumerate(table.rows):
        for c, cell in enumerate(row.cells):
            if cell.is_spanned:
                continue  # 被合并覆盖的格不产 cell
            if cell.is_merge_origin:
                row_span, col_span = int(cell.span_height), int(cell.span_width)
            else:
                row_span, col_span = 1, 1
            cells.append({
                "row_index": r,
                "column_index": c,
                "text": cell.text.strip(),
                "row_span": row_span,
                "column_span": col_span,
                "is_header": r == 0,  # 首行表头约定
                "evidence_index": evidence,
            })
            evidence += 1
    structure: dict[str, Any] = {"rows": rows, "cols": cols, "cells": cells}
    return BackendBlock(
        block_type="table",
        text="",
        structure=structure,
        **common,
    )


def _bbox_of(shape: BaseShape) -> tuple[float, float, float, float] | None:
    """shape 位置 -> ``(x0, top, x1, bottom)`` 角点（整改轮修复 I-1）.

    无位置的 shape 返回 None 不伪造。
    """
    if None in (shape.left, shape.top, shape.width, shape.height):
        return None
    left, top = float(shape.left), float(shape.top)
    return (left, top, left + float(shape.width), top + float(shape.height))


def _is_title_placeholder(shape: BaseShape) -> bool:
    if not shape.is_placeholder:
        return False
    return shape.placeholder_format.type in _TITLE_PLACEHOLDER_TYPES


# ---------------------------------------------------------------------------
# Normalizer（SRS §C07）
# ---------------------------------------------------------------------------

class PptxNormalizer(BaseNativeNormalizer):
    """PPTX backend artifact -> Parse IR：slide 容器 + EMU 坐标证据."""

    normalizer_version = "native-pptx@2"
    _default_fingerprints = {NATIVE_PPTX_PARSER_ID: NATIVE_PPTX_FINGERPRINT}
    # slide 是独立内容单元：父链不跨 slide（避免无标题 slide 的正文
    # 误挂上一张 slide 的标题，评审 MED）。
    reset_heading_stack_on_container_change = True
    _element_type_map = {
        "heading": "heading",
        "paragraph": "paragraph",
        "list_item": "list_item",
        "table": "table",
        "figure": "figure",
    }

    def _build_containers(self, artifact) -> tuple[Container, ...]:
        slide_indices = sorted({
            (b.container_ref or {}).get("index")
            for b in artifact.blocks if b.container_ref
        }, key=lambda v: (v is None, v))
        slide_ids = {index: f"c-slide-{index}" for index in slide_indices}
        width = artifact.usage.get("slide_width_emu")
        height = artifact.usage.get("slide_height_emu")
        return tuple(
            Container(
                container_id=slide_ids[index],
                container_type="slide",
                order_index=int(index),
                name=None,  # 不伪造页码/名称（SRS §3.6）
                width=float(width) if width is not None else None,
                height=float(height) if height is not None else None,
                coordinate_unit="emu",
            )
            for index in slide_indices
        )

    def _container_id_for(self, block, containers) -> str | None:
        index = (block.container_ref or {}).get("index")
        for container in containers:
            if container.container_type == "slide":
                if container.order_index == index:
                    return container.container_id
        return None

    def _make_spans(self, element_id, block, container_id) -> tuple[EvidenceSpan, ...]:
        visual = None
        if block.bbox is not None:
            visual = {
                "bbox": [block.bbox[0], block.bbox[1], block.bbox[2], block.bbox[3]],
                "unit": "emu",
            }
        return (EvidenceSpan(
            span_id=f"{element_id}-s0",
            page_id=container_id,
            visual_region=visual,
            native_ref=dict(block.native_ref) if block.native_ref else None,
            raw_text=block.text or None,
        ),)

    def _make_cell_spans(
        self, element_id, block, container_id
    ) -> tuple[EvidenceSpan, ...]:
        """表格 cell 级 EvidenceSpan（slide/shape + row/col 定位，I-4）."""
        if block.block_type != "table":
            return ()
        structure = block.structure or {}
        raw_cells = structure.get("cells") or []
        native = dict(block.native_ref or {})
        spans: list[EvidenceSpan] = []
        for k, cell in enumerate(raw_cells):
            ref = dict(native)
            ref["row_index"] = int(cell["row_index"])
            ref["column_index"] = int(cell["column_index"])
            spans.append(EvidenceSpan(
                span_id=f"{element_id}-cell-{k:04d}",
                native_ref=ref,
                raw_text=str(cell.get("text", "")) or None,
            ))
        return tuple(spans)

    def _extra_assets(self, element_id, block, container_id) -> dict[str, Any]:
        """figure 块 -> FigureAsset + binary_assets 条目（整改轮）."""
        if block.block_type != "figure":
            return {}
        structure = block.structure or {}
        sha = structure.get("image_sha256")
        if not sha:
            return {}
        binary_id = structure.get("binary_asset_id", f"bin-{sha[:16]}")
        asset = FigureAsset(
            figure_id=f"{element_id}-figure",
            original_image_ref=binary_id,
            image_hash=sha,
        )
        return {"asset": asset, "binary": (binary_id, {
            "sha256": sha,
            "mime": structure.get("mime"),
            "size": structure.get("size"),
            "origin": NATIVE_PPTX_PARSER_ID,
        })}


__all__ = [
    "NATIVE_PPTX_FINGERPRINT",
    "NATIVE_PPTX_MIMES",
    "NATIVE_PPTX_PARSER_ID",
    "NATIVE_PPTX_VERSION",
    "NativePptxParser",
    "PptxNormalizer",
]
