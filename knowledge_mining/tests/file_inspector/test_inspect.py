"""Unit tests for FileInspector（M3 WP5, SRS §C03 / §4.2 裁剪字段）.

纯逻辑测试，所有 fixture 在测试内于内存生成（BytesIO，不落盘）：

- OOXML 用 python-docx / openpyxl / python-pptx 现场构造；
- PDF 用手写最小合法 PDF 字节串（带 xref 表 + 文本内容流，构造说明
  见 ``_build_pdf`` 注释；pdfminer 只要求 xref 偏移正确即可解析）。

覆盖：各格式 source_format/container 探测、加密 PDF 标记、伪扩展名被
签名识破（SRS §2.4 "扩展名不可信"）、未知格式不抛（SRS §4.4 unsupported
语义）、无 declared_mime 时 OOXML 由 [Content_Types].xml 消歧。
"""
from __future__ import annotations

import io

import docx as python_docx
import openpyxl
import pptx
import pytest

from knowledge_mining.mining.file_inspector.inspect import (
    DocumentProfile,
    FileInspector,
)

OOXML_DOCX_MIME = (
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
)
OOXML_XLSX_MIME = (
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
)
OOXML_PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


# ---------------------------------------------------------------------------
# 内存 fixture 构造（bytes，不落盘）
# ---------------------------------------------------------------------------


def _build_docx_bytes() -> bytes:
    document = python_docx.Document()
    document.add_paragraph("第一段")
    document.add_paragraph("第二段")
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _build_xlsx_bytes(sheet_count: int = 2) -> bytes:
    workbook = openpyxl.Workbook()
    while len(workbook.sheetnames) < sheet_count:
        workbook.create_sheet()
    buf = io.BytesIO()
    workbook.save(buf)
    return buf.getvalue()


def _build_pptx_bytes(slide_count: int = 2) -> bytes:
    presentation = pptx.Presentation()
    while len(presentation.slides._sldIdLst) < slide_count:  # noqa: SLF001
        presentation.slides.add_slide(
            presentation.slide_layouts[6]  # blank layout
        )
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


def _build_pdf(page_count: int = 2, with_text: bool = True) -> bytes:
    """手写最小合法 PDF 字节串（不依赖外部 fixture 文件）。

    构成：1=Pages 树, 2..N+1=Page, N+2..2N+1=内容流, 2N+2=Font；带文本
    时内容流为 ``BT /F1 .. Tf (..) Tj ET``——pdfminer 抽字符**必须**有
    /Font 资源（无字体时 chars 恒为空，这一点 test_pdf_parser_images.py
    的 fixture 同样遵守）；无文本时只画一个矩形（``re``+``f``）。最后按
    对象偏移生成 xref 表 + trailer——pdfminer 解析只要求 xref 偏移准确。
    """
    objects: list[bytes] = []
    page_ids = list(range(2, 2 + page_count))
    kids = " ".join(f"{pid} 0 R" for pid in page_ids)
    objects.append(
        b"<< /Type /Pages /Kids [" + kids.encode() + b"] /Count "
        + str(page_count).encode() + b" >>"
    )
    font_id = 2 + 2 * page_count
    stream_ids = list(range(2 + page_count, font_id))
    for pid, sid in zip(page_ids, stream_ids):
        objects.append(
            b"<< /Type /Page /Parent 1 0 R /MediaBox [0 0 200 100] "
            b"/Resources << /Font << /F1 "
            + f"{font_id} 0 R".encode()
            + b" >> >> /Contents " + f"{sid} 0 R".encode() + b" >>"
        )
    for _ in range(page_count):
        if with_text:
            content = b"BT /F1 12 Tf 10 50 Td (hello pdf) Tj ET\n"
        else:
            content = b"10 10 100 5 re f\n"
        objects.append(
            f"<< /Length {len(content)} >>\nstream\n".encode()
            + content + b"endstream"
        )
    objects.append(
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>"
    )
    out = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for i, obj in enumerate(objects, start=1):
        offsets.append(len(out))
        out.extend(f"{i} 0 obj\n".encode())
        out.extend(obj)
        out.extend(b"\nendobj\n")
    xref_pos = len(out)
    out.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    out.extend(b"0000000000 65535 f \n")
    for off in offsets[1:]:
        out.extend(f"{off:010d} 00000 n \n".encode())
    out.extend(
        f"trailer<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref_pos}\n%%EOF\n".encode()
    )
    return bytes(out)


@pytest.fixture
def inspector() -> FileInspector:
    return FileInspector()


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def test_pdf_with_text_layer(inspector: FileInspector) -> None:
    profile = inspector.inspect(_build_pdf(page_count=2, with_text=True))
    assert profile.detected_mime == "application/pdf"
    assert profile.source_format == "pdf"
    assert profile.container_kind == "page"
    assert profile.container_count == 2
    assert profile.encrypted is False
    assert profile.has_text_layer is True


def test_pdf_without_text_layer(inspector: FileInspector) -> None:
    profile = inspector.inspect(_build_pdf(page_count=1, with_text=False))
    assert profile.source_format == "pdf"
    assert profile.has_text_layer is False


def test_encrypted_pdf_flagged(inspector: FileInspector) -> None:
    # 在 trailer 里注入 /Encrypt 引用（指向不存在的对象也无妨：探测在
    # pdfplumber 打开前完成，与 safe_intake 的 admission 启发一致）
    raw = _build_pdf(page_count=1)
    data = raw.replace(b"trailer<<", b"trailer<</Encrypt 99 0 R")
    profile = inspector.inspect(data)
    assert profile.encrypted is True
    assert profile.source_format == "pdf"


# ---------------------------------------------------------------------------
# OOXML
# ---------------------------------------------------------------------------


def test_docx_profile(inspector: FileInspector) -> None:
    profile = inspector.inspect(_build_docx_bytes(), declared_mime=OOXML_DOCX_MIME)
    assert profile.detected_mime == OOXML_DOCX_MIME
    assert profile.source_format == "docx"
    # 段落不是容器（SRS §4.2：docx 无原生容器语义）
    assert profile.container_count is None
    assert profile.container_kind is None
    assert profile.encrypted is False
    assert profile.has_text_layer is None


def test_xlsx_sheet_count(inspector: FileInspector) -> None:
    profile = inspector.inspect(_build_xlsx_bytes(sheet_count=3),
                                declared_mime=OOXML_XLSX_MIME)
    assert profile.source_format == "xlsx"
    assert profile.container_kind == "sheet"
    assert profile.container_count == 3


def test_pptx_slide_count(inspector: FileInspector) -> None:
    profile = inspector.inspect(_build_pptx_bytes(slide_count=2),
                                declared_mime=OOXML_PPTX_MIME)
    assert profile.source_format == "pptx"
    assert profile.container_kind == "slide"
    assert profile.container_count == 2


def test_ooxml_detected_without_declared_mime(inspector: FileInspector) -> None:
    # ZIP 签名 + [Content_Types].xml 消歧：即使调用方未声明 MIME 也能识别
    for data, expected in (
        (_build_docx_bytes(), "docx"),
        (_build_xlsx_bytes(), "xlsx"),
        (_build_pptx_bytes(), "pptx"),
    ):
        profile = inspector.inspect(data)
        assert profile.source_format == expected, expected


def test_corrupt_ooxml_zip_is_unknown_not_crash(inspector: FileInspector) -> None:
    # PK 头但 central directory 损坏：不抛，降级为 unknown + warning
    data = b"PK\x03\x04" + b"\x00" * 64
    profile = inspector.inspect(data)
    assert profile.source_format == "unknown"
    assert profile.warnings


# ---------------------------------------------------------------------------
# 文本格式
# ---------------------------------------------------------------------------


def test_text_formats_have_no_container(inspector: FileInspector) -> None:
    cases = [
        (b"<html><body><p>hi</p></body></html>", "text/html", "html"),
        (b"# Title\n\nbody", "text/markdown", "markdown"),
        (b"plain text line", "text/plain", "txt"),
    ]
    for data, mime, fmt in cases:
        profile = inspector.inspect(data, declared_mime=mime)
        assert profile.source_format == fmt
        assert profile.container_count is None
        assert profile.container_kind is None
        assert profile.has_text_layer is None


def test_spoofed_extension_exposed_by_signature(inspector: FileInspector) -> None:
    # 声明是 PDF，实为文本：签名/文本启发胜出（SRS §2.4 扩展名不可信）
    profile = inspector.inspect(b"# not a pdf\njust text", declared_mime="application/pdf")
    assert profile.source_format != "pdf"
    assert profile.source_format in {"txt", "markdown"}


def test_unknown_format_does_not_raise(inspector: FileInspector) -> None:
    # 含 NUL 的二进制：非文本、无签名、无可用扩展提示 → unknown，不抛
    profile = inspector.inspect(b"\x00\x01\x02\x03\xde\xad\xbe\xef" * 8)
    assert profile.source_format == "unknown"
    assert profile.container_count is None
    assert profile.container_kind is None


def test_profile_is_frozen(inspector: FileInspector) -> None:
    profile = inspector.inspect(b"abc", declared_mime="text/plain")
    assert isinstance(profile, DocumentProfile)
    with pytest.raises(Exception):
        profile.source_format = "pdf"  # type: ignore[misc]
