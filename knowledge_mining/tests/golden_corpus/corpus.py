"""Golden corpus（整改轮，用户指令第 5 步）：50 份多类型确定性文档.

每格式含**正例 / 反例 / 复杂结构 / 退化输入**四类；全部在内存确定性
构造（无网络、无随机、无外部文件），期望标注（标题序列/段落锚文本/
表格数）随文档返回，供 benchmark 计算结构准确率。

- MD/TXT/HTML：字节字面量；
- DOCX/XLSX/PPTX：python-docx / openpyxl / python-pptx 构造；
- PDF：最小 PDF 构造器（与 test_native_pdf 相同模板，Helvetica——
  CJK 无法在 fixture 中渲染，PDF 样本为英文形态）。
"""
from __future__ import annotations

import io
from dataclasses import dataclass
from typing import Any

from knowledge_mining.mining.parse_quality import GoldenExpectations


@dataclass(frozen=True)
class CorpusDoc:
    """一份 golden 文档：字节 + 期望 + 类别标注."""

    name: str
    format_key: str  # md/txt/docx/xlsx/pptx/html/pdf
    mime: str
    data: bytes
    expectations: GoldenExpectations
    category: str  # positive / negative / complex / degenerate
    source_text: str | None = None  # 字符覆盖率对照（文本可提取格式）


MIME = {
    "md": "text/markdown",
    "txt": "text/plain",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "html": "text/html",
    "pdf": "application/pdf",
}

PARSER_ID = {
    "md": "legacy_markdown",
    "txt": "legacy_txt",
    "docx": "native_docx",
    "xlsx": "native_xlsx",
    "pptx": "native_pptx",
    "html": "native_html",
    "pdf": "native_pdf",
}


# ---------------------------------------------------------------------------
# MD（7）
# ---------------------------------------------------------------------------


def _md_docs() -> list[CorpusDoc]:
    def doc(name, text, exp, category, anchors=()):
        return CorpusDoc(
            name, "md", MIME["md"], text.encode("utf-8"),
            GoldenExpectations(
                expected_headings=exp, expected_paragraph_anchors=anchors,
            ),
            category, source_text=text,
        )

    return [
        doc(
            "md-headings-lists",
            "# 概述\n\n引导段落。\n\n## 安装\n\n- 步骤一\n- 步骤二\n\n"
            "## 配置\n\n说明文字。\n",
            ("概述", "安装", "配置"), "positive", ("引导段落", "步骤一"),
        ),
        doc(
            "md-table",
            "# 数据\n\n| 名称 | 值 |\n| --- | --- |\n| 甲 | 1 |\n| 乙 | 2 |\n",
            ("数据",), "positive", ("甲", "乙"),
        ),
        CorpusDoc(
            "md-code-quote", "md", MIME["md"],
            "# 杂项\n\n> 引用一段。\n\n```python\nprint(1)\n```\n".encode(
                "utf-8"
            ),
            GoldenExpectations(
                expected_headings=("杂项",),
                expected_paragraph_anchors=("引用一段", "print"),
            ),
            "positive",
            # 覆盖率不适用：code fence 语言标注不进元素文本
        ),
        CorpusDoc(
            "md-image-link-diag", "md", MIME["md"],
            "![图](a.png)[链](https://x.y)\n".encode("utf-8"),
            GoldenExpectations(), "complex",
            # 覆盖率不适用：inline 图片/链接的标记语法字符非实义内容
        ),
        doc(
            "md-nested-list",
            "# 列表\n\n- 甲\n  - 甲一\n  - 甲二\n- 乙\n",
            ("列表",), "complex", ("甲一", "乙"),
        ),
        CorpusDoc(
            "md-empty", "md", MIME["md"], b"", GoldenExpectations(), "degenerate",
        ),
        CorpusDoc(
            "md-garbage", "md", MIME["md"], b"\xff\xfe\x00binary-ish",
            GoldenExpectations(), "negative",
        ),
    ]


# ---------------------------------------------------------------------------
# TXT（5）
# ---------------------------------------------------------------------------


def _txt_docs() -> list[CorpusDoc]:
    return [
        CorpusDoc(
            "txt-paragraphs", "txt", MIME["txt"],
            "第一段内容。\n\n第二段内容。\n\n第三段内容。\n".encode("utf-8"),
            GoldenExpectations(expected_paragraph_anchors=("第一段", "第三段")),
            "positive", source_text="第一段内容。第二段内容。第三段内容。",
        ),
        CorpusDoc(
            "txt-empty", "txt", MIME["txt"], b"", GoldenExpectations(),
            "degenerate",
        ),
        CorpusDoc(
            "txt-single-long-line", "txt", MIME["txt"],
            ("超长单行段落，" * 400).encode("utf-8"),
            GoldenExpectations(expected_paragraph_anchors=("超长单行",)),
            "complex",
        ),
        CorpusDoc(
            "txt-blank-only", "txt", MIME["txt"], b"\n\n\n  \n",
            GoldenExpectations(), "degenerate",
        ),
        CorpusDoc(
            "txt-crlf", "txt", MIME["txt"],
            "Windows 段落一。\r\n\r\nWindows 段落二。\r\n".encode("utf-8"),
            GoldenExpectations(expected_paragraph_anchors=("段落一", "段落二")),
            "positive", source_text="Windows 段落一。Windows 段落二。",
        ),
    ]


# ---------------------------------------------------------------------------
# DOCX（8）
# ---------------------------------------------------------------------------


def _docx_bytes(build) -> bytes:
    from docx import Document

    doc = Document()
    build(doc)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _docx_docs() -> list[CorpusDoc]:
    def mk(name, build, exp, category, anchors=()):
        return CorpusDoc(
            name, "docx", MIME["docx"], _docx_bytes(build),
            GoldenExpectations(
                expected_headings=exp, expected_paragraph_anchors=anchors,
            ),
            category,
        )

    def heading_tree(d):
        d.add_heading("章一", level=1)
        d.add_paragraph("章一导语。")
        d.add_heading("节1.1", level=2)
        d.add_paragraph("节内容。")

    def lists(d):
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn

        d.add_paragraph("编号列表前导。")
        for i, (text, ilvl) in enumerate(
            [("顶层甲", 0), ("顶层乙", 0), ("嵌套项", 1)]
        ):
            p = d.add_paragraph(text)
            pPr = p._p.get_or_add_pPr()
            numPr = OxmlElement("w:numPr")
            ilvl_el = OxmlElement("w:ilvl")
            ilvl_el.set(qn("w:val"), str(ilvl))
            numid = OxmlElement("w:numId")
            numid.set(qn("w:val"), "1")
            numPr.append(ilvl_el)
            numPr.append(numid)
            pPr.append(numPr)

    def merged_table(d):
        d.add_paragraph("合并表格说明。")
        t = d.add_table(rows=3, cols=3)
        t.cell(0, 0).merge(t.cell(0, 1))
        t.cell(0, 0).text = "表头横跨"
        t.cell(0, 2).text = "C"
        t.cell(1, 0).text = "a"
        t.cell(1, 1).text = "b"
        t.cell(1, 2).text = "c"
        t.cell(2, 0).text = "d"
        t.cell(2, 1).text = "e"
        t.cell(2, 2).text = "f"

    def nested_table(d):
        d.add_paragraph("嵌套表宿主。")
        outer = d.add_table(rows=2, cols=2)
        outer.cell(0, 0).text = "外A"
        outer.cell(0, 1).text = "外B"
        outer.cell(1, 0).text = "宿主"
        outer.cell(1, 1).text = "外D"
        inner = outer.cell(1, 0).add_table(rows=1, cols=2)
        inner.cell(0, 0).text = "内1"
        inner.cell(0, 1).text = "内2"

    def with_image(d):
        from PIL import Image

        d.add_paragraph("图片文档正文。")
        img = io.BytesIO()
        Image.new("RGB", (2, 2), (9, 9, 9)).save(img, format="PNG")
        d.add_picture(io.BytesIO(img.getvalue()))

    def with_header_footer(d):
        d.add_paragraph("正文一段。")
        d.sections[0].header.paragraphs[0].text = "页眉文字内容"
        d.sections[0].footer.paragraphs[0].text = "页脚文字内容"

    def complex_doc(d):
        d.add_heading("复杂文档", level=1)
        d.add_paragraph("复杂导语段。")
        t = d.add_table(rows=2, cols=2)
        t.cell(0, 0).text = "列甲"
        t.cell(0, 1).text = "列乙"
        t.cell(1, 0).text = "1"
        t.cell(1, 1).text = "2"
        d.add_heading("次级标题", level=2)
        d.add_paragraph("次级内容。")

    def empty(d):
        pass

    return [
        mk("docx-heading-tree", heading_tree, ("章一", "节1.1"), "positive",
           ("导语", "节内容")),
        mk("docx-numpr-lists", lists, (), "positive", ("顶层甲", "嵌套项")),
        mk("docx-merged-table", merged_table, (), "positive",
           ("表头横跨", "合并表格")),
        mk("docx-nested-table", nested_table, (), "complex", ("外A", "内1")),
        mk("docx-image-diag", with_image, (), "complex", ("图片文档",)),
        mk("docx-header-footer-diag", with_header_footer, (), "complex",
           ("正文一段",)),
        mk("docx-complex", complex_doc, ("复杂文档", "次级标题"), "complex",
           ("复杂导语", "次级内容")),
        mk("docx-empty", empty, (), "degenerate"),
    ]


# ---------------------------------------------------------------------------
# XLSX（7）
# ---------------------------------------------------------------------------


def _xlsx_bytes(build) -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    build(wb)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _xlsx_docs() -> list[CorpusDoc]:
    def mk(name, build, tables, category, anchors=()):
        return CorpusDoc(
            name, "xlsx", MIME["xlsx"], _xlsx_bytes(build),
            GoldenExpectations(
                expected_table_count=tables, expected_paragraph_anchors=anchors,
            ),
            category,
        )

    def single_region(wb):
        ws = wb.active
        ws.title = "数据"
        for row in [("名称", "值"), ("甲", 1), ("乙", 2)]:
            ws.append(row)

    def multi_region(wb):
        ws = wb.active
        ws["A1"] = "表一头"
        ws["B1"] = "表一值"
        ws["A2"] = "x"
        ws["B2"] = "1"
        ws["D5"] = "表二头"
        ws["E5"] = "表二值"
        ws["D6"] = "y"
        ws["E6"] = "2"

    def excel_table(wb):
        from openpyxl.worksheet.table import Table

        ws = wb.active
        ws["A1"] = "列甲"
        ws["B1"] = "列乙"
        ws["A2"] = "1"
        ws["B2"] = "2"
        ws["E9"] = "表外杂散"
        ws.add_table(Table(displayName="命名表", ref="A1:B2"))

    def formulas(wb):
        ws = wb.active
        ws["A1"] = "数量"
        ws["B1"] = "合计"
        ws["A2"] = 3
        ws["B2"] = "=A2*2"

    def merged_hidden(wb):
        ws = wb.active
        ws["A1"] = "锚点"
        ws.merge_cells("A1:B2")
        ws["C3"] = "明文"
        ws["C4"] = "隐藏行"
        ws.row_dimensions[4].hidden = True

    def with_chart(wb):
        from openpyxl.chart import BarChart, Reference

        ws = wb.active
        ws["A1"] = "k"
        ws["B1"] = "v"
        ws["A2"] = "a"
        ws["B2"] = 1
        chart = BarChart()
        chart.add_data(Reference(ws, min_col=2, min_row=1, max_row=2))
        ws.add_chart(chart, "D2")

    def sparse_bomb(wb):
        ws = wb.active
        ws["A1"] = "近"
        ws.cell(row=1, column=16384, value="远")

    return [
        mk("xlsx-single-region", single_region, 1, "positive", ("名称",)),
        mk("xlsx-multi-region", multi_region, 2, "positive", ("表一头", "表二头")),
        mk("xlsx-excel-table", excel_table, 1, "positive", ("列甲",)),
        mk("xlsx-formulas", formulas, 1, "complex", ("合计",)),
        mk("xlsx-merged-hidden", merged_hidden, 1, "complex", ("锚点", "明文")),
        mk("xlsx-chart-diag", with_chart, 1, "complex", ("k",)),
        mk("xlsx-sparse-bomb", sparse_bomb, 2, "negative"),
    ]


# ---------------------------------------------------------------------------
# PPTX（7）
# ---------------------------------------------------------------------------


def _pptx_bytes(build) -> bytes:
    from pptx import Presentation

    prs = Presentation()
    build(prs)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _pptx_docs() -> list[CorpusDoc]:
    def mk(name, build, exp, category, anchors=()):
        return CorpusDoc(
            name, "pptx", MIME["pptx"], _pptx_bytes(build),
            GoldenExpectations(
                expected_headings=exp, expected_paragraph_anchors=anchors,
            ),
            category,
        )

    def title_body(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "首页标题"
        box = slide.shapes.add_textbox(914400, 914400, 4000000, 900000)
        box.text_frame.text = "正文框内容"

    def bullets(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(100000, 100000, 6000000, 4000000)
        tf = box.text_frame
        tf.text = "引导行"
        p2 = tf.add_paragraph()
        p2.text = "一级要点"
        p2.level = 1
        p3 = tf.add_paragraph()
        p3.text = "二级要点"
        p3.level = 2

    def table(prs):
        from pptx.util import Emu

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        t = slide.shapes.add_table(
            2, 2, Emu(100000), Emu(100000), Emu(4000000), Emu(1800000)
        ).table
        t.cell(0, 0).text = "头甲"
        t.cell(0, 1).text = "头乙"
        t.cell(1, 0).text = "1"
        t.cell(1, 1).text = "2"

    def picture(prs):
        from PIL import Image
        from pptx.util import Emu

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img = io.BytesIO()
        Image.new("RGB", (4, 4), (1, 2, 3)).save(img, format="PNG")
        slide.shapes.add_picture(
            io.BytesIO(img.getvalue()), Emu(1), Emu(1), Emu(914400), Emu(914400)
        )

    def notes(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(1, 1, 4000000, 900000)
        box.text_frame.text = "可见内容"
        slide.notes_slide.notes_text_frame.text = "讲者备注内容"

    def group_shapes(prs):
        from pptx.util import Emu

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        group = slide.shapes.add_group_shape()
        box = group.shapes.add_textbox(
            Emu(100000), Emu(100000), Emu(2000000), Emu(500000)
        )
        box.text_frame.text = "组内文本"

    def chart_slide(prs):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Emu

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        caption = slide.shapes.add_textbox(1, 1, 4000000, 500000)
        caption.text_frame.text = "图表页说明"
        cd = CategoryChartData()
        cd.categories = ["a", "b"]
        cd.add_series("s", (1, 2))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Emu(1), Emu(1), Emu(3000000), Emu(2000000), cd,
        )

    return [
        mk("pptx-title-body", title_body, ("首页标题",), "positive", ("正文框",)),
        mk("pptx-bullets", bullets, (), "positive", ("引导行", "二级要点")),
        mk("pptx-table", table, (), "positive", ("头甲",)),
        mk("pptx-picture", picture, (), "complex", ()),
        mk("pptx-notes", notes, (), "complex", ("可见内容", "讲者备注")),
        mk("pptx-group", group_shapes, (), "complex", ("组内文本",)),
        mk("pptx-chart-diag", chart_slide, (), "complex", ()),
    ]


# ---------------------------------------------------------------------------
# HTML（8）
# ---------------------------------------------------------------------------


def _html_docs() -> list[CorpusDoc]:
    def doc(name, html, exp, category, anchors=(), tables=None):
        return CorpusDoc(
            name, "html", MIME["html"], html.encode("utf-8"),
            GoldenExpectations(
                expected_headings=exp, expected_paragraph_anchors=anchors,
                expected_table_count=tables,
            ),
            category, source_text=None,
        )

    return [
        doc(
            "html-basic",
            "<html><head><title>页面</title></head><body><h1>主题</h1>"
            "<p>正文段落。</p></body></html>",
            ("主题",), "positive", ("正文段落",),
        ),
        doc(
            "html-nested-list",
            "<html><body><ul><li>水果<ul><li>苹果</li><li>香蕉</li></ul></li>"
            "<li>蔬菜</li></ul></body></html>",
            (), "complex", ("苹果", "蔬菜"),
        ),
        doc(
            "html-table-merge",
            "<html><body><table><tr><th>甲</th><th>乙</th></tr>"
            "<tr><td rowspan='2'>跨两行</td><td>1</td></tr>"
            "<tr><td>2</td></tr></table></body></html>",
            (), "positive", ("跨两行",), tables=1,
        ),
        doc(
            "html-figure-caption",
            "<html><body><figure><img src='f.png' alt='图'>"
            "<figcaption>图注文字</figcaption></figure>"
            "<table><caption>表题文字</caption>"
            "<tr><th>a</th><th>b</th></tr><tr><td>1</td><td>2</td></tr>"
            "</table></body></html>",
            (), "complex", ("图注文字", "表题文字"), tables=1,
        ),
        doc(
            "html-links",
            "<html><body><h1>链接页</h1>"
            "<p>参见 <a href='https://e.x/d'>文档</a> 与 "
            "<a href='https://e.x/s'>规范</a>。</p></body></html>",
            ("链接页",), "positive", ("参见",),
        ),
        doc(
            "html-semantic",
            "<html><body><article><section><h2>章节</h2>"
            "<p>语义段落。</p></section></article></body></html>",
            ("章节",), "complex", ("语义段落",),
        ),
        CorpusDoc(
            "html-rowspan-bomb", "html", MIME["html"],
            ("<html><body><table><tr><td rowspan='9999' colspan='9999'>大</td>"
             "</tr><tr><td>x</td></tr></table></body></html>").encode("utf-8"),
            GoldenExpectations(expected_table_count=1), "negative",
        ),
        CorpusDoc(
            "html-empty", "html", MIME["html"],
            b"<html><body></body></html>", GoldenExpectations(), "degenerate",
        ),
    ]


# ---------------------------------------------------------------------------
# PDF（8）：最小 PDF 构造器（Helvetica / 英文形态）
# ---------------------------------------------------------------------------


def _build_pdf(objects: list[bytes]) -> bytes:
    out = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for i, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n".encode() + body + b"\nendobj\n"
    xref_pos = len(out)
    n = len(objects) + 1
    out += f"xref\n0 {n}\n".encode()
    out += b"0000000000 65535 f \n"
    for off in offsets[1:]:
        out += f"{off:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {n} /Root 1 0 R >>\n"
        f"startxref\n{xref_pos}\n%%EOF\n"
    ).encode()
    return bytes(out)


def _stream(content: str) -> bytes:
    return f"<< /Length {len(content)} >>\nstream\n{content}\nendstream".encode()


def _t(size: float, x: float, y: float, text: str) -> str:
    escaped = text.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")
    return f"BT /F1 {size} Tf {x} {y} Td ({escaped}) Tj ET"


def _page(content: str) -> bytes:
    return _build_pdf([
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
         b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        _stream(content),
    ])


def _two_page_pdf(contents: list[str]) -> bytes:
    kids = " ".join(f"{3 + 2 * i} 0 R" for i in range(len(contents)))
    objects: list[bytes] = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        f"<< /Type /Pages /Kids [{kids}] /Count {len(contents)} >>".encode(),
    ]
    # 对象编号：1 catalog, 2 pages, (page, stream) 对，最后 font。
    # page i 的 Contents = 其 stream 对象 = 4 + 2*i；字体对象放最后
    # （编号 3 + 2*n），页面对象经资源名 /F1 引用。
    font_num = 3 + 2 * len(contents)
    for i, content in enumerate(contents):
        objects.append(
            (f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
             f"/Resources << /Font << /F1 {font_num} 0 R >> >> /Contents "
             f"{4 + 2 * i} 0 R >>").encode()
        )
        objects.append(_stream(content))
    objects.append(
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>"
    )
    return _build_pdf(objects)


def _pdf_docs() -> list[CorpusDoc]:
    def doc(name, data, exp, category, anchors=(), tables=None):
        return CorpusDoc(
            name, "pdf", MIME["pdf"], data,
            GoldenExpectations(
                expected_headings=exp, expected_paragraph_anchors=anchors,
                expected_table_count=tables,
            ),
            category,
        )

    single_col = _page("\n".join([
        _t(16, 72, 700, "Chapter One Overview"),
        _t(10, 72, 660, "Body first line of the single column page."),
        _t(10, 72, 645, "Body second line continues the same paragraph."),
    ]))
    two_col = _page("\n".join([
        _t(14, 72, 700, "Two Column Layout"),
        _t(10, 72, 670, "L1 left column first"),
        _t(10, 72, 650, "L2 left column second"),
        _t(10, 72, 630, "L3 left column third"),
        _t(10, 72, 610, "L4 left column fourth"),
        _t(10, 350, 670, "R1 right column one"),
        _t(10, 350, 650, "R2 right column two"),
        _t(10, 350, 630, "R3 right column three"),
        _t(10, 350, 610, "R4 right column four"),
    ]))
    line_table = _page("\n".join([
        _t(10, 80, 665, "HeadA"), _t(10, 180, 665, "HeadB"),
        _t(10, 80, 635, "a"), _t(10, 180, 635, "b"),
        _t(10, 80, 605, "c"), _t(10, 180, 605, "d"),
        "0 w",
        "72 600 m 272 600 l S", "72 633 m 272 633 l S",
        "72 667 m 272 667 l S", "72 700 m 272 700 l S",
        "72 600 m 72 700 l S", "172 600 m 172 700 l S",
        "272 600 m 272 700 l S",
    ]))
    trap = _page("\n".join([
        _t(12, 80, 690, "Trap Title"),
        _t(10, 80, 660, "HeadA"), _t(10, 180, 660, "HeadB"),
        _t(10, 80, 630, "a"), _t(10, 180, 630, "b"),
        _t(10, 80, 600, "c"), _t(10, 180, 600, "d"),
        "0 w",
        "72 590 m 272 590 l S", "72 700 m 272 700 l S",
    ]))
    captioned = _page("\n".join([
        _t(10, 100, 700, "Table 2-1 Sample data summary"),
        _t(10, 80, 665, "HeadA"), _t(10, 180, 665, "HeadB"),
        _t(10, 80, 635, "a"), _t(10, 180, 635, "b"),
        _t(10, 80, 605, "c"), _t(10, 180, 605, "d"),
        "0 w",
        "72 600 m 272 600 l S", "72 633 m 272 633 l S",
        "72 667 m 272 667 l S", "72 700 m 272 700 l S",
        "72 600 m 72 700 l S", "172 600 m 172 700 l S",
        "272 600 m 272 700 l S",
    ]))
    furniture = _two_page_pdf([
        "\n".join([
            _t(10, 200, 40, "Shared header line text"),
            _t(10, 72, 660, "First page body content line."),
        ]),
        "\n".join([
            _t(10, 200, 40, "Shared header line text"),
            _t(10, 72, 660, "Second page body content line."),
        ]),
    ])
    numbered = _page("\n".join([
        _t(10, 72, 700, "1.1 Background of study"),
        _t(10, 72, 670, "2.1 Methods and materials"),
        _t(10, 72, 640, "Body line to establish modal size."),
        _t(10, 72, 625, "Another body line for the modal baseline."),
    ]))
    blank = _build_pdf([
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R 5 0 R] /Count 2 >>",
        (b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
         b"/Resources << /Font << /F1 4 0 R >> >> /Contents 6 0 R >>"),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        (b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
         b"/Resources << >> /Contents 7 0 R >>"),
        _stream(_t(10, 72, 700, "Only text page.")),
        _stream(""),
    ])

    return [
        doc("pdf-single-col", single_col, ("Chapter One Overview",), "positive",
            ("single column",)),
        doc("pdf-two-col", two_col, ("Two Column Layout",), "complex",
            ("left column", "right column")),
        doc("pdf-line-table", line_table, (), "positive", ("HeadA",), tables=1),
        doc("pdf-trap-table", trap, (), "negative", ("HeadA",), tables=1),
        doc("pdf-captioned-table", captioned, (), "complex", ("Sample data",),
            tables=1),
        doc("pdf-furniture", furniture, (), "complex", ("body content",)),
        doc("pdf-numbered-headings", numbered,
            ("1.1 Background", "2.1 Methods"), "complex", ("modal size",)),
        doc("pdf-blank-page", blank, (), "degenerate"),
    ]


# ---------------------------------------------------------------------------
# 汇总
# ---------------------------------------------------------------------------


def build_corpus() -> list[CorpusDoc]:
    """全量 golden corpus（确定性：同代码同输出）。"""
    return [
        *_md_docs(), *_txt_docs(), *_docx_docs(), *_xlsx_docs(),
        *_pptx_docs(), *_html_docs(), *_pdf_docs(),
    ]


def corpus_stats() -> dict[str, Any]:
    docs = build_corpus()
    by_format: dict[str, int] = {}
    by_category: dict[str, int] = {}
    for d in docs:
        by_format[d.format_key] = by_format.get(d.format_key, 0) + 1
        by_category[d.category] = by_category.get(d.category, 0) + 1
    return {
        "total": len(docs),
        "by_format": by_format,
        "by_category": by_category,
    }


__all__ = ["CorpusDoc", "build_corpus", "corpus_stats", "PARSER_ID", "MIME"]
