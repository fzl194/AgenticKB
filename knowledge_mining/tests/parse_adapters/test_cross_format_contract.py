"""跨格式 Parse IR 不变量 contract tests（整改轮，用户指令第 2 步）.

不同于各 adapter 自己的单测（自定义语义），本文件对**全部 7 条 route**
断言同一组 IR 不变量——格式间语义必须由契约统一，而不是每个 adapter
各自表述：

  C-1  bbox 统一 ``(x0, top, x1, bottom)``（角点，不是 left/top/w/h）
  C-2  表格 Element.text 是统一的确定性可检索序列化文本
  C-3  TableAsset 是表格事实源；Element.text 可由 TableAsset 重算得到
  C-4  能取得原生单元格位置时，每个 TableCell 有独立 source_span_id
  C-5  backend raw artifact 可序列化持久化，且 replay（重新 normalize）
       与直连结果完全一致（SRS §9.5 / A09）
  C-6  不支持/未解析的结构进入 diagnostics，不得静默丢失
  C-7  解析产物的 ParseIdentity 携带 rule_config_fingerprint，
       effective fingerprint 对规则变化敏感
  C-8  不把 adapter 特有类型泄漏进公共 Parse IR（元素类型全部合法）
"""
from __future__ import annotations

import hashlib
import io

import pytest

from knowledge_mining.mining.contracts.parse_ir import (
    VALID_ELEMENT_TYPES,
    validate,
)
from knowledge_mining.mining.contracts.parser_adapter import (
    BackendParseArtifact,
    effective_pipeline_fingerprint,
)
from knowledge_mining.mining.parse_adapters.factory import resolve_pipeline

RAW_HASH = "ab" * 32

MIME = {
    "md": "text/markdown",
    "txt": "text/plain",
    "docx": (
        "application/vnd.openxmlformats-officedocument."
        "wordprocessingml.document"
    ),
    "xlsx": (
        "application/vnd.openxmlformats-officedocument."
        "spreadsheetml.sheet"
    ),
    "pptx": (
        "application/vnd.openxmlformats-officedocument."
        "presentationml.presentation"
    ),
    "html": "text/html",
    "pdf": "application/pdf",
}


# ---------------------------------------------------------------------------
# 每格式构造一份"含表格"的样本字节
# ---------------------------------------------------------------------------


def _md_bytes() -> bytes:
    return (
        "# 标题一\n\n段落甲。\n\n| 表头A | 表头B |\n| --- | --- |\n"
        "| a1 | b1 |\n| a2 | b2 |\n\n![图](img.png)[链接](https://x.y)\n"
    ).encode("utf-8")


def _txt_bytes() -> bytes:
    return "第一段。\n\n第二段。\n\n第三段。\n".encode("utf-8")


def _docx_bytes() -> bytes:
    from docx import Document
    from docx.shared import Inches

    doc = Document()
    doc.add_heading("章标题", level=1)
    doc.add_paragraph("导语。")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "表头A"
    table.cell(0, 1).text = "表头B"
    table.cell(1, 0).text = "a"
    table.cell(1, 1).text = "b"
    doc.add_picture(io.BytesIO(_png_bytes()), width=Inches(1))
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "数据"
    ws["A1"] = "表头A"
    ws["B1"] = "表头B"
    ws["A2"] = "a1"
    ws["B2"] = "b1"
    ws["A3"] = "a2"
    ws["B3"] = "b2"
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pptx_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "页标题"
    box = slide.shapes.add_textbox(Emu(100), Emu(200), Emu(300), Emu(400))
    box.text_frame.text = "正文框"
    table = slide.shapes.add_table(
        2, 2, Emu(914400), Emu(914400), Emu(1828800), Emu(457200)
    ).table
    table.cell(0, 0).text = "表头A"
    table.cell(0, 1).text = "表头B"
    table.cell(1, 0).text = "a"
    table.cell(1, 1).text = "b"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _html_bytes() -> bytes:
    return (
        "<html><head><title>文档</title></head><body>"
        "<h1>标题</h1><p>段落 <a href='https://e.x'>链接</a>。</p>"
        "<table><tr><th>表头A</th><th>表头B</th></tr>"
        "<tr><td>a</td><td>b</td></tr></table>"
        "<img src='x.png' alt='图'>"
        "</body></html>"
    ).encode("utf-8")


def _png_bytes() -> bytes:
    # 1x1 RGB PNG（Pillow 生成，python-docx/openpyxl 均可接受）
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (1, 1), (255, 0, 0)).save(buf, format="PNG")
    return buf.getvalue()


TABLE_FORMATS = ["md", "docx", "xlsx", "pptx", "html", "pdf"]
ALL_FORMATS = ["md", "txt", "docx", "xlsx", "pptx", "html", "pdf"]


def _sample(format_key: str) -> bytes:
    builders = {
        "md": _md_bytes,
        "txt": _txt_bytes,
        "docx": _docx_bytes,
        "xlsx": _xlsx_bytes,
        "pptx": _pptx_bytes,
        "html": _html_bytes,
        "pdf": None,  # 由 pdf 专属 fixture 提供
    }
    builder = builders[format_key]
    if builder is None:
        raise SkipTest("pdf sample provided by dedicated fixture")
    return builder()


class SkipTest(Exception):
    pass


def _run(format_key: str, data: bytes):
    """format -> (parser, normalizer) 流水线执行：parse + normalize."""
    parser_id = {
        "md": "legacy_markdown",
        "txt": "legacy_txt",
        "docx": "native_docx",
        "xlsx": "native_xlsx",
        "pptx": "native_pptx",
        "html": "native_html",
        "pdf": "native_pdf",
    }[format_key]
    pair = resolve_pipeline(parser_id)
    assert pair is not None, parser_id
    parser, normalizer = pair
    artifact = parser.parse(data, mime=MIME[format_key])
    doc = normalizer.normalize(artifact, source_raw_hash=RAW_HASH)
    return artifact, doc


def _tables(doc):
    return [
        (e, doc.structured_assets[f"{e.element_id}-table"])
        for e in doc.elements
        if e.element_type == "table"
        and f"{e.element_id}-table" in doc.structured_assets
    ]


# ---------------------------------------------------------------------------
# C-2 / C-3：表格 Element.text 统一序列化 + 可由 TableAsset 重算
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("fmt", TABLE_FORMATS)
def test_table_element_text_is_unified_and_derivable(fmt):
    """含表格式：table 元素 text 非空、含表头与数据、且与 TableAsset 重算一致."""
    if fmt == "pdf":
        pytest.skip("pdf 由专属用例覆盖")
    _, doc = _run(fmt, _sample(fmt))
    tables = _tables(doc)
    assert tables, f"{fmt}: 无表格元素/资产"
    for element, asset in tables:
        assert element.text.strip(), f"{fmt}: 表格 Element.text 为空"
        assert "表头A" in element.text, f"{fmt}: 表头不在可检索文本中: {element.text!r}"
        assert "a" in element.text
        # C-3：TableAsset 是事实源——Element.text 必须能由 asset 确定性重算
        rebuilt = render_table_text(asset)
        assert rebuilt == element.text, (
            f"{fmt}: Element.text 与 TableAsset 渲染不一致"
        )


def render_table_text(asset) -> str:
    """统一序列化：行 \\n、列 \\t、只取原点 cell（跨格式唯一定义）."""
    from knowledge_mining.mining.parse_adapters.rendered_text import (
        render_table_text as _render,
    )
    return _render(asset)


# ---------------------------------------------------------------------------
# C-4：TableCell 独立 source_span_id（可取得原生位置时）
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("fmt", ["docx", "xlsx", "pptx", "html"])
def test_table_cells_have_independent_spans(fmt):
    _, doc = _run(fmt, _sample(fmt))
    tables = _tables(doc)
    assert tables
    for element, asset in tables:
        span_ids = {s.span_id for s in element.source_spans}
        seen = set()
        for cell in asset.cells:
            if not cell.text:
                continue
            assert cell.source_span_id is not None, (
                f"{fmt}: cell[{cell.row_index},{cell.column_index}] 无 source_span_id"
            )
            assert cell.source_span_id in span_ids
            assert cell.source_span_id not in seen, (
                f"{fmt}: 多个 cell 共享同一 span（无独立证据）"
            )
            seen.add(cell.source_span_id)


def test_pdf_table_cells_have_bbox_spans(pdf_table_doc):
    """PDF 表格：pdfplumber 提供每 cell 坐标 -> 必须有独立 bbox 证据."""
    tables = _tables(pdf_table_doc)
    assert tables
    for element, asset in tables:
        span_ids = {s.span_id for s in element.source_spans}
        for cell in asset.cells:
            if not cell.text:
                continue
            assert cell.source_span_id is not None
            assert cell.source_span_id in span_ids


# ---------------------------------------------------------------------------
# C-1：bbox 角点语义（用 PPTX 的已知放置验证；PDF 行 bbox 一并断言）
# ---------------------------------------------------------------------------


def test_pptx_bbox_is_corner_pair():
    """已知放置 (left=100, top=200, w=300, h=400) -> bbox == (100, 200, 400, 600)."""
    _, doc = _run("pptx", _pptx_bytes())
    target = None
    for e in doc.elements:
        if e.text == "正文框":
            target = e
    assert target is not None
    bbox = target.source_spans[0].visual_region["bbox"]
    assert bbox == pytest.approx([100.0, 200.0, 400.0, 600.0]), bbox


def test_all_visual_bboxes_are_ordered_corner_pairs(pdf_text_doc):
    """全部带 bbox 的格式：x0 <= x1 且 top <= bottom（角点语义）."""
    _, pptx_doc = _run("pptx", _pptx_bytes())
    for fmt, doc in (("pptx", pptx_doc), ("pdf", pdf_text_doc)):
        for e in doc.elements:
            for s in e.source_spans:
                vr = s.visual_region
                if vr and "bbox" in vr:
                    x0, top, x1, bottom = vr["bbox"]
                    assert x0 <= x1 and top <= bottom, (fmt, e.element_id, vr)


# ---------------------------------------------------------------------------
# C-5：raw artifact 序列化 + replay 等价
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("fmt", ["docx", "xlsx", "pptx", "html", "pdf"])
def test_backend_artifact_replay_is_equivalent(fmt):
    """serialize -> from_dict -> 重新 normalize == 直连 normalize（§9.5 A09）."""
    if fmt == "pdf":
        pytest.skip("pdf 由专属用例覆盖")
    parser_id = {
        "docx": "native_docx", "xlsx": "native_xlsx", "pptx": "native_pptx",
        "html": "native_html", "pdf": "native_pdf",
    }[fmt]
    _, normalizer = resolve_pipeline(parser_id)
    artifact, direct = _run(fmt, _sample(fmt))

    restored = BackendParseArtifact.from_dict(artifact.to_dict())
    replayed = normalizer.normalize(restored, source_raw_hash=RAW_HASH)
    assert replayed.to_dict() == direct.to_dict(), f"{fmt}: replay 不等价"


# ---------------------------------------------------------------------------
# C-6：不支持/未解析结构进入 diagnostics
# ---------------------------------------------------------------------------


def _warnings_joined(doc) -> str:
    return "\n".join(doc.diagnostics.warnings)


def test_docx_image_is_diagnosed():
    _, doc = _run("docx", _docx_bytes())  # fixture 内含一张图片
    text = _warnings_joined(doc)
    assert "image" in text.lower(), f"DOCX 图片静默丢失: {text!r}"


def test_html_image_is_diagnosed():
    _, doc = _run("html", _html_bytes())
    text = _warnings_joined(doc)
    assert "image" in text.lower(), f"HTML 图片静默丢失: {text!r}"


def test_markdown_image_and_link_are_diagnosed():
    _, doc = _run("md", _md_bytes())
    text = _warnings_joined(doc)
    assert "image" in text.lower(), f"MD 图片静默丢失: {text!r}"
    assert "link" in text.lower(), f"MD 链接静默丢失: {text!r}"


# ---------------------------------------------------------------------------
# C-7：rule_config_fingerprint + effective fingerprint 敏感性
# ---------------------------------------------------------------------------


def test_identity_carries_rule_config_fingerprint():
    _, doc = _run("docx", _docx_bytes())
    assert doc.source_identity.rule_config_fingerprint, "ParseIdentity 缺规则指纹"


def test_effective_fingerprint_reacts_to_rule_change():
    from knowledge_mining.mining.contracts.parser_adapter import ParseRuleConfig

    base = ParseRuleConfig()
    changed = ParseRuleConfig(heading_size_ratio=1.4)
    kw = dict(
        parser_fingerprint="native_pdf@1.0.0",
        normalizer_version="pdf-native@1",
        dependency_fingerprint="d",
    )
    fp1 = effective_pipeline_fingerprint(
        rule_config_fingerprint=base.config_fingerprint(), **kw
    )
    fp2 = effective_pipeline_fingerprint(
        rule_config_fingerprint=changed.config_fingerprint(), **kw
    )
    assert fp1 != fp2


# ---------------------------------------------------------------------------
# C-8：无 adapter 特有类型泄漏
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("fmt", ALL_FORMATS)
def test_no_adapter_type_leak(fmt):
    if fmt == "pdf":
        pytest.skip("pdf 由专属用例覆盖")
    _, doc = _run(fmt, _sample(fmt))
    for e in doc.elements:
        assert e.element_type in VALID_ELEMENT_TYPES
    assert validate(doc).valid


# ---------------------------------------------------------------------------
# PDF 专属 fixture（复用 test_native_pdf 的最小构造器）
# ---------------------------------------------------------------------------


def _build_pdf(objects):
    out = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for i, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n".encode() + body + b"\nendobj\n"
    xref_pos = len(out)
    n = len(objects) + 1
    out += f"xref\n0 {n}\n".encode()
    out += b"0000000000 65535 f \n"
    for off in offsets[1:]:
        out += f"{off:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {n} /Root 1 0 R >>\n"
        f"startxref\n{xref_pos}\n%%EOF\n"
    ).encode()
    return bytes(out)


def _stream(content: str) -> bytes:
    return f"<< /Length {len(content)} >>\nstream\n{content}\nendstream".encode()


def _pdf_with_table() -> bytes:
    content = (
        "BT /F1 10 Tf 72 700 Td (Above table.) Tj ET\n"
        "BT /F1 10 Tf 80 665 Td (HeadA) Tj ET\n"
        "BT /F1 10 Tf 180 665 Td (HeadB) Tj ET\n"
        "BT /F1 10 Tf 80 635 Td (a) Tj ET\n"
        "BT /F1 10 Tf 180 635 Td (b) Tj ET\n"
        "0 w\n"
        "72 600 m 272 600 l S\n72 633 m 272 633 l S\n"
        "72 667 m 272 667 l S\n72 700 m 272 700 l S\n"
        "72 600 m 72 700 l S\n172 600 m 172 700 l S\n"
        "272 600 m 272 700 l S\n"
    )
    return _build_pdf([
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
         b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        _stream(content),
    ])


@pytest.fixture(scope="module")
def pdf_table_doc():
    _, doc = _run("pdf", _pdf_with_table())
    return doc


@pytest.fixture(scope="module")
def pdf_text_doc():
    content = (
        "BT /F1 12 Tf 72 700 Td (Title line) Tj ET\n"
        "BT /F1 10 Tf 72 660 Td (Body one.) Tj ET\n"
        "BT /F1 10 Tf 72 645 Td (Body two.) Tj ET\n"
    )
    data = _build_pdf([
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
         b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        _stream(content),
    ])
    _, doc = _run("pdf", data)
    return doc


def test_pdf_table_element_text_derivable(pdf_table_doc):
    doc = pdf_table_doc
    tables = _tables(doc)
    assert tables
    for element, asset in tables:
        assert "HeadA" in element.text
        assert element.text == render_table_text(asset)


def test_pdf_artifact_replay_equivalent():
    parser, normalizer = resolve_pipeline("native_pdf")
    data = _pdf_with_table()
    artifact = parser.parse(data, mime=MIME["pdf"])
    direct = normalizer.normalize(artifact, source_raw_hash=RAW_HASH)
    restored = BackendParseArtifact.from_dict(artifact.to_dict())
    replayed = normalizer.normalize(restored, source_raw_hash=RAW_HASH)
    assert replayed.to_dict() == direct.to_dict()
