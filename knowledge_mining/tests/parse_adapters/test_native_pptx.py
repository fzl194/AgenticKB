"""``NativePptxParser`` + ``PptxNormalizer`` 单元测试（M3, SRS §C06/§C07）.

fixture 用 python-pptx 在测试内生成 bytes。覆盖：slide 容器（无页码伪造）、
title placeholder -> heading level 1、文本 shape -> paragraph、
bbox 保留 EMU 坐标、表格合并、IR 断言链、错误归一。
"""
from __future__ import annotations

import io

import pytest

from knowledge_mining.mining.contracts.parse_ir import validate
from knowledge_mining.mining.contracts.parser_adapter import (
    ParserAdapterError,
    UnsupportedFormat,
)
from knowledge_mining.mining.parse_adapters.native.native_pptx import (
    NATIVE_PPTX_FINGERPRINT,
    NativePptxParser,
    PptxNormalizer,
)

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
RAW_HASH = "99" * 32

TITLE_L, TITLE_T, TITLE_W, TITLE_H = 100, 200, 300, 400  # EMU
BOX_L, BOX_T, BOX_W, BOX_H = 914400, 914400, 1828800, 457200  # EMU
BOX2_L, BOX2_T, BOX2_W, BOX2_H = 10, 20, 500, 600  # EMU


def _build_pptx_bytes() -> bytes:
    """2 页：页一 title+文本框+2x2 合并表格；页二（blank）文本框."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "幻灯一标题"
    box = slide1.shapes.add_textbox(
        Emu(BOX_L), Emu(BOX_T), Emu(BOX_W), Emu(BOX_H)
    )
    box.text_frame.text = "第一页正文"
    table = slide1.shapes.add_table(
        2, 2, Emu(TITLE_L), Emu(TITLE_T), Emu(TITLE_W), Emu(TITLE_H)
    ).table
    table.cell(0, 0).merge(table.cell(0, 1))
    table.cell(0, 0).text = "表头"
    table.cell(1, 0).text = "a"
    table.cell(1, 1).text = "b"

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box2 = slide2.shapes.add_textbox(
        Emu(BOX2_L), Emu(BOX2_T), Emu(BOX2_W), Emu(BOX2_H)
    )
    box2.text_frame.text = "第二页文本"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def pptx_bytes() -> bytes:
    return _build_pptx_bytes()


@pytest.fixture
def parser() -> NativePptxParser:
    return NativePptxParser()


@pytest.fixture
def normalizer() -> PptxNormalizer:
    return PptxNormalizer()


def _normalize(parser, normalizer, data: bytes):
    artifact = parser.parse(data, mime=PPTX_MIME)
    return normalizer.normalize(artifact, source_raw_hash=RAW_HASH)


# ---------------------------------------------------------------------------
# RED 1: parse 层
# ---------------------------------------------------------------------------

def test_parse_blocks_in_geometric_reading_order(parser, pptx_bytes) -> None:
    """整改轮：阅读序 = 几何带排序，不等于 shape XML 顺序.

    fixture 中表格位于页顶（top=200 EMU）-> 几何序第一；title 占位符
    默认 top≈274638 第二；正文框 top=914400 第三；第二页第四。
    """
    artifact = parser.parse(pptx_bytes, mime=PPTX_MIME)

    assert artifact.parser_id == "native_pptx"
    # slide2 无标题占位符：页角(top=20 EMU)短文本按 PPT 惯例升 heading（v9）
    assert [b.block_type for b in artifact.blocks] == [
        "table", "heading", "paragraph", "heading",
    ]
    assert artifact.blocks[0].container_ref == {"container_type": "slide", "index": 0}
    assert artifact.blocks[3].container_ref == {"container_type": "slide", "index": 1}
    assert artifact.blocks[1].level == 1
    # bbox 角点语义：(x0, top, x1, bottom) = (left, top, left+w, top+h)
    assert artifact.blocks[2].bbox == (
        float(BOX_L), float(BOX_T), float(BOX_L + BOX_W), float(BOX_T + BOX_H),
    )
    # native_ref 保留 XML shape_index（回溯原始结构）
    assert artifact.blocks[2].native_ref == {"slide_index": 0, "shape_index": 2}


def test_parse_unsupported_mime_and_corrupt(parser) -> None:
    with pytest.raises(UnsupportedFormat):
        parser.parse(b"x", mime="application/pdf")
    with pytest.raises(ParserAdapterError):
        parser.parse(b"garbage-not-ooxml", mime=PPTX_MIME)


# ---------------------------------------------------------------------------
# RED 2: IR 断言链
# ---------------------------------------------------------------------------

def test_slide_containers_no_fake_page_numbers(parser, normalizer, pptx_bytes) -> None:
    doc = _normalize(parser, normalizer, pptx_bytes)
    assert validate(doc).valid

    assert [c.container_type for c in doc.containers] == ["slide", "slide"]
    for i, container in enumerate(doc.containers):
        assert container.page_number is None  # 不伪造页码
        assert container.order_index == i
        assert container.coordinate_unit == "emu"
    # 幻灯片物理尺寸（EMU）记录在容器上
    assert doc.containers[0].width == 9144000
    assert doc.containers[0].height == 6858000


def test_element_types_heading_paragraph_table(parser, normalizer, pptx_bytes) -> None:
    doc = _normalize(parser, normalizer, pptx_bytes)

    assert [e.element_type for e in doc.elements] == [
        "table", "heading", "paragraph", "heading",
    ]
    heading = doc.elements[1]
    assert heading.text == "幻灯一标题"
    assert heading.style["level"] == 1
    assert heading.parent_id is None
    assert doc.elements[2].text == "第一页正文"
    assert doc.elements[3].text == "第二页文本"  # heading（v9 页角惯例）


def test_bbox_emu_evidence_and_slide_pages(parser, normalizer, pptx_bytes) -> None:
    doc = _normalize(parser, normalizer, pptx_bytes)

    body = doc.elements[2]
    span = body.source_spans[0]
    assert span.visual_region == {
        "bbox": [BOX_L, BOX_T, BOX_L + BOX_W, BOX_T + BOX_H],
        "unit": "emu",
    }
    # 元素挂 slide 容器：span.page_id 与 page_span_ids 都指向 slide 容器
    assert span.page_id == doc.containers[0].container_id
    assert body.page_span_ids == (doc.containers[0].container_id,)
    body2 = doc.elements[3]
    assert body2.page_span_ids == (doc.containers[1].container_id,)


def test_table_asset_with_merge(parser, normalizer, pptx_bytes) -> None:
    doc = _normalize(parser, normalizer, pptx_bytes)

    table_el = doc.elements[0]
    asset = doc.structured_assets[f"{table_el.element_id}-table"]
    assert asset.rows == 2
    assert asset.columns == 2
    grid = {(c.row_index, c.column_index): c for c in asset.cells}
    assert grid[(0, 0)].text == "表头"
    assert grid[(0, 0)].column_span == 2
    assert (0, 1) not in grid
    assert grid[(1, 0)].text == "a"
    assert grid[(1, 1)].text == "b"
    assert grid[(0, 0)].is_header


def test_ir_chain_relations_and_fingerprint(parser, normalizer, pptx_bytes) -> None:
    doc = _normalize(parser, normalizer, pptx_bytes)
    assert validate(doc).valid

    reading = [
        r for r in doc.relations if r.relation_type == "next_in_reading_order"
    ]
    assert len(reading) == 3
    ids = [e.element_id for e in doc.elements]
    assert [(r.source_element_id, r.target_element_id) for r in reading] == list(
        zip(ids, ids[1:])
    )
    assert doc.source_identity.parser_fingerprint == NATIVE_PPTX_FINGERPRINT
    assert doc.source_identity.parser_fingerprint == (
        "native_pptx@2.0.0#python-pptx-1.0.2"
    )


def test_top_zone_short_line_promoted_heading() -> None:
    """页面顶部短行升级 heading（PPT 排版惯例，验收 v9）.

    真实答辩 PPT 的标题是普通文本框（非占位符）——通用信号：
    位于页面上部 30% 区域、单行、显著短（<=24 字符）。
    """
    from knowledge_mining.mining.parse_adapters.native.native_pptx import (
        _looks_like_slide_title,
    )

    # 顶部短行（top=0.08×页高）-> heading
    assert _looks_like_slide_title("研究背景", top_emu=720000, slide_h=9000000) is True
    # 顶部长行 -> 不是
    assert _looks_like_slide_title(
        "增量学习目的在于开发人工智能系统可以不断地从新数据中学习",
        top_emu=720000, slide_h=9000000,
    ) is False
    # 底部短行（页码区）-> 不是
    assert _looks_like_slide_title("3", top_emu=8500000, slide_h=9000000) is False


# ===========================================================================
# 整改轮（2026-08-17）：bbox 角点 / 段落拆分 / 阅读序 / notes / group / 图片资产
# ===========================================================================


def _prs():
    from pptx import Presentation

    return Presentation()


def test_text_frame_paragraphs_split_with_bullet_levels() -> None:
    """文本框内多段落拆分；bullet 段落 -> list_item（层级来自 a:pPr/lvl）."""
    from pptx.util import Emu

    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(100000), Emu(200000), Emu(5000000), Emu(3000000))
    tf = box.text_frame
    tf.text = "引导段落"
    p2 = tf.add_paragraph()
    p2.text = "要点一"
    p2.level = 1
    p3 = tf.add_paragraph()
    p3.text = "子要点"
    p3.level = 2
    buf = io.BytesIO(); prs.save(buf)

    artifact = NativePptxParser().parse(buf.getvalue(), mime=PPTX_MIME)
    texts_levels = [(b.text, b.block_type, b.level) for b in artifact.blocks]
    assert ("引导段落", "paragraph", None) in texts_levels
    assert ("要点一", "list_item", 2) in texts_levels
    assert ("子要点", "list_item", 3) in texts_levels


def test_reading_order_is_geometric_not_xml_order() -> None:
    """XML 顺序 = 右框先写；几何阅读序必须上带先、带内左先."""
    from pptx.util import Emu

    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 先加"下一行左框"，再加"上一行右框"，再加"上一行左框"
    bottom_right = slide.shapes.add_textbox(Emu(4000000), Emu(3000000), Emu(900000), Emu(400000))
    bottom_right.text_frame.text = "下右"
    top_right = slide.shapes.add_textbox(Emu(4000000), Emu(100000), Emu(900000), Emu(400000))
    top_right.text_frame.text = "上右"
    top_left = slide.shapes.add_textbox(Emu(100000), Emu(100000), Emu(900000), Emu(400000))
    top_left.text_frame.text = "上左"
    buf = io.BytesIO(); prs.save(buf)

    artifact = NativePptxParser().parse(buf.getvalue(), mime=PPTX_MIME)
    texts = [b.text for b in artifact.blocks if b.text]
    assert texts == ["上左", "上右", "下右"], f"阅读序错误: {texts}"


def test_slide_notes_are_preserved() -> None:
    from pptx.util import Emu

    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(1), Emu(1), Emu(900000), Emu(400000))
    box.text_frame.text = "正文"
    slide.notes_slide.notes_text_frame.text = "这是演讲备注"
    buf = io.BytesIO(); prs.save(buf)

    artifact = NativePptxParser().parse(buf.getvalue(), mime=PPTX_MIME)
    notes_blocks = [
        b for b in artifact.blocks
        if (b.structure or {}).get("slide_notes") is True
    ]
    assert notes_blocks and notes_blocks[0].text == "这是演讲备注"


def test_group_shapes_are_recursed() -> None:
    from pptx.util import Emu

    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    box = group.shapes.add_textbox(Emu(100000), Emu(100000), Emu(900000), Emu(400000))
    box.text_frame.text = "组内文本"
    buf = io.BytesIO(); prs.save(buf)

    artifact = NativePptxParser().parse(buf.getvalue(), mime=PPTX_MIME)
    texts = [b.text for b in artifact.blocks]
    assert "组内文本" in texts, f"group 成员丢失: {texts}"


def test_picture_produces_figure_asset_with_image_hash() -> None:
    from pptx.util import Emu
    from PIL import Image

    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img = io.BytesIO()
    Image.new("RGB", (3, 3), (10, 20, 30)).save(img, format="PNG")
    slide.shapes.add_picture(
        io.BytesIO(img.getvalue()), Emu(100000), Emu(100000), Emu(914400), Emu(914400)
    )
    buf = io.BytesIO(); prs.save(buf)

    import hashlib
    expected_hash = hashlib.sha256(img.getvalue()).hexdigest()

    parser = NativePptxParser()
    artifact = parser.parse(buf.getvalue(), mime=PPTX_MIME)
    doc = PptxNormalizer().normalize(artifact, source_raw_hash="99" * 32)
    figures = [e for e in doc.elements if e.element_type == "figure"]
    assert figures, "图片元素缺失"
    figure_assets = [
        a for a in doc.structured_assets.values()
        if getattr(a, "figure_id", "").endswith("-figure")
    ]
    assert figure_assets, "FigureAsset 缺失"
    asset = figure_assets[0]
    assert asset.image_hash == expected_hash
    assert expected_hash[:16] in "".join(doc.binary_assets.keys())


def test_chart_is_diagnosed_not_silent() -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Emu

    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = ["a", "b"]
    cd.add_series("s", (1, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Emu(1), Emu(1), Emu(900000), Emu(400000), cd
    )
    buf = io.BytesIO(); prs.save(buf)

    artifact = NativePptxParser().parse(buf.getvalue(), mime=PPTX_MIME)
    joined = "\n".join(artifact.warnings).lower()
    assert "chart" in joined, f"图表静默丢失: {joined!r}"


def test_table_cells_have_independent_spans() -> None:
    parser = NativePptxParser()
    doc = PptxNormalizer().normalize(
        parser.parse(_build_pptx_bytes(), mime=PPTX_MIME), source_raw_hash="99" * 32
    )
    table = next(e for e in doc.elements if e.element_type == "table")
    asset = doc.structured_assets[f"{table.element_id}-table"]
    span_by_id = {s.span_id: s for s in table.source_spans}
    seen = set()
    for cell in asset.cells:
        if not cell.text:
            continue
        assert cell.source_span_id is not None
        span = span_by_id[cell.source_span_id]
        assert "row_index" in span.native_ref
        assert cell.source_span_id not in seen
        seen.add(cell.source_span_id)
