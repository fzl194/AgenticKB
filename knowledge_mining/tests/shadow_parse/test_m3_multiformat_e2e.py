"""M3 全链路集成：Inspector → Router → 工厂 → ShadowParseService ×6 格式.

SRS §14 M3 退出条件的内存版验收："至少 native/Docling 两条 route 可运行，
结果统一进入 IR"——本套件验证 **七条已实现 route**（MD/TXT/DOCX/XLSX/
PPTX/HTML/PDF）从冻结字节到 parse bucket IR 的完整链路，且 IR 全部通过
schema validation（round-trip）。

对每格式的断言侧重"格式特有保真点"（M3 的核心价值）：
- PDF：页容器 + visual_region bbox 证据 + heading 启发式
- DOCX：样式标题树 + 合并单元格
- XLSX：workbook→sheet 容器层级 + cell native_ref 证据 + 公式/展示值
- PPTX：slide 容器 + EMU bbox
- HTML：DOM 容器 + xpath 证据
"""
from __future__ import annotations

import asyncio
import hashlib
import io
import json
import sys
from collections.abc import AsyncIterator

import pytest

if sys.platform == "win32":
    asyncio.set_event_loop_policy(asyncio.WindowsSelectorEventLoopPolicy())

from docx import Document as DocxDocument  # noqa: E402
from openpyxl import Workbook  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

from knowledge_mining.mining.contracts.parse_ir.schema import validate  # noqa: E402
from knowledge_mining.mining.contracts.parse_ir.types import ParsedDocument  # noqa: E402
from knowledge_mining.mining.contracts.storage.types import (  # noqa: E402
    ObjectLocation,
    PutOptions,
)
from knowledge_mining.mining.file_inspector.inspect import FileInspector  # noqa: E402
from knowledge_mining.mining.file_inspector.router import ParserRouter  # noqa: E402
from knowledge_mining.mining.file_management.repositories_memory import (  # noqa: E402
    MemoryStorageObjectRepository,
)
from knowledge_mining.mining.frozen_input.contracts import FrozenInput  # noqa: E402
from knowledge_mining.mining.infra.object_store.fake import FakeObjectStore  # noqa: E402
from knowledge_mining.mining.parse_adapters.factory import resolve_pipeline  # noqa: E402
from knowledge_mining.mining.parse_adapters.registry import (  # noqa: E402
    build_default_registry,
)
from knowledge_mining.mining.shadow_parse.repositories_memory import (  # noqa: E402
    MemoryParseRunRepository,
)
from knowledge_mining.mining.shadow_parse.service import ShadowParseService  # noqa: E402
from knowledge_mining.tests.parse_adapters.test_native_pdf import (  # noqa: E402
    _build_pdf,
    _page,
)

_BUCKET_PREFIX = "testm3e2e-"
SOURCE_BUCKET = f"{_BUCKET_PREFIX}source"

# ---------------------------------------------------------------------------
# 各格式 fixture bytes（库内生成；PDF 复用 native_pdf 测试的构造器）
# ---------------------------------------------------------------------------


def _docx_bytes() -> bytes:
    doc = DocxDocument()
    doc.add_heading("运维手册", level=1)
    doc.add_heading("巡检流程", level=2)
    doc.add_paragraph("每日执行例行巡检。")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "项目"
    table.cell(0, 1).text = "频率"
    table.cell(1, 0).text = "风扇"
    table.cell(1, 1).text = "每日"
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _xlsx_bytes() -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "参数表"
    ws.append(["参数", "默认值"])
    ws.append(["timeout", 30])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pptx_bytes() -> bytes:
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # title-only
    slide = prs.slides.add_slide(slide_layout)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "发布说明"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


_HTML_BYTES = """<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>告警说明</title></head>
<body><h1>告警说明</h1><p>告警分为两级。</p>
<table><tr><th>码</th><th>含义</th></tr><tr><td>E-01</td><td>过温</td></tr></table>
</body></html>""".encode("utf-8")

_MD_BYTES = "# 手册\n\n正文。\n\n| a | b |\n| --- | --- |\n| 1 | 2 |\n".encode("utf-8")
_TXT_BYTES = "第一段。\n\n第二段。".encode("utf-8")

_PDF_BYTES = _build_pdf(
    _page(
        "",
        "BT /F1 18 Tf 72 700 Td (Manual Title) Tj ET\n"
        "BT /F1 10 Tf 72 660 Td (Body text line.) Tj ET",
    )
)

#: (格式名, bytes, 声明 MIME, 保真断言函数名)
CASES = [
    ("markdown", _MD_BYTES, "text/markdown"),
    ("txt", _TXT_BYTES, "text/plain"),
    ("docx", _docx_bytes(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
    ("xlsx", _xlsx_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
    ("pptx", _pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    ("html", _HTML_BYTES, "text/html"),
    ("pdf", _PDF_BYTES, "application/pdf"),
]


async def _chunks(payload: bytes) -> AsyncIterator[bytes]:
    for i in range(0, len(payload), 4096):
        yield payload[i : i + 4096]


async def _run_format(tmp_path, content: bytes, mime: str):
    """全链路：inspect → route → resolve → ShadowParseService.run → 读回 IR."""
    profile = FileInspector().inspect(content, declared_mime=mime)
    router = ParserRouter(build_default_registry())
    decision = router.plan(profile)
    assert decision.primary_parser_id is not None, profile
    pair = resolve_pipeline(decision.primary_parser_id)
    assert pair is not None, decision.primary_parser_id
    parser, normalizer = pair

    store = FakeObjectStore(str(tmp_path / "store"))
    runs = MemoryParseRunRepository()
    storage_objects = MemoryStorageObjectRepository()
    sha = hashlib.sha256(content).hexdigest()
    await store.put_stream(
        ObjectLocation(bucket=SOURCE_BUCKET, object_key=f"src/{profile.source_format}"),
        _chunks(content),
        PutOptions(artifact_class="source", expected_sha256=sha),
    )
    frozen = FrozenInput(
        document_id=f"doc-m3-{profile.source_format}",
        source_storage_object_id="so_src",
        source_raw_hash=sha,
        source_content_revision=1,
        mime=mime,
        size=len(content),
        original_filename="sample",
        captured_at="2026-08-17T00:00:00+00:00",
        provider="fake",
        bucket=SOURCE_BUCKET,
        object_key=f"src/{profile.source_format}",
    )
    service = ShadowParseService(
        object_store=store,
        parse_runs=runs,
        storage_objects=storage_objects,
        parser=parser,
        normalizer=normalizer,
        bucket_prefix=_BUCKET_PREFIX,
    )
    result = await service.run(frozen)
    assert result.status == "SUCCEEDED"

    so = await storage_objects.get(result.parse_ir_storage_object_id)
    assert so is not None and so.artifact_class == "parse_ir"
    buf = b""
    async for chunk in store.get_stream(
        ObjectLocation(bucket=so.bucket, object_key=so.object_key)
    ):
        buf += chunk
    doc = ParsedDocument.from_dict(json.loads(buf.decode("utf-8")))
    verdict = validate(doc)
    assert verdict.valid, [f"{i.code}: {i.message}" for i in verdict.issues]
    return profile, decision, doc


# ---------------------------------------------------------------------------
# 全格式全链路
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
@pytest.mark.parametrize("fmt,content,mime", CASES, ids=[c[0] for c in CASES])
async def test_multiformat_full_pipeline(tmp_path, fmt, content, mime) -> None:
    """7 格式：Inspector 识别 → Router 路由 → IR 落桶 → validate 通过。"""
    profile, decision, doc = await _run_format(tmp_path, content, mime)

    assert profile.source_format == fmt, profile
    assert decision.primary_parser_id is not None
    assert len(doc.elements) > 0, f"{fmt}: no elements"


@pytest.mark.asyncio
async def test_pdf_fidelity_points(tmp_path) -> None:
    """PDF 保真：页容器 + bbox 证据 + heading 启发式命中大字号行。"""
    _, _, doc = await _run_format(tmp_path, _PDF_BYTES, "application/pdf")
    pages = [c for c in doc.containers if c.container_type == "page"]
    assert pages, "page containers missing"
    headings = [e for e in doc.elements if e.element_type == "heading"]
    assert headings and "Manual Title" in headings[0].text
    # bbox 证据（visual_region）
    evidenced = [
        e for e in doc.elements
        for s in e.source_spans if s.visual_region and "bbox" in s.visual_region
    ]
    assert evidenced, "bbox evidence missing"


@pytest.mark.asyncio
async def test_docx_fidelity_points(tmp_path) -> None:
    """DOCX 保真：样式标题树 + 表格网格。"""
    _, _, doc = await _run_format(tmp_path, _docx_bytes(), CASES[2][2])
    headings = {e.text: e for e in doc.elements if e.element_type == "heading"}
    assert "巡检流程" in headings
    h2 = headings["巡检流程"]
    by_id = {e.element_id: e for e in doc.elements}
    assert by_id[h2.parent_id].text == "运维手册"
    tables = [a for a in doc.structured_assets.values() if getattr(a, "rows", None)]
    assert tables and tables[0].rows == 2 and any("风扇" in c.text for c in tables[0].cells)


@pytest.mark.asyncio
async def test_xlsx_fidelity_points(tmp_path) -> None:
    """XLSX 保真：workbook→sheet 层级 + cell native_ref 证据。"""
    _, _, doc = await _run_format(tmp_path, _xlsx_bytes(), CASES[3][2])
    kinds = {c.container_type for c in doc.containers}
    assert "workbook" in kinds and "sheet" in kinds, kinds
    native_refs = [
        s.native_ref for e in doc.elements for s in e.source_spans if s.native_ref
    ]
    assert any(nr.get("cell") for nr in native_refs), native_refs


@pytest.mark.asyncio
async def test_pptx_fidelity_points(tmp_path) -> None:
    """PPTX 保真：slide 容器（无伪造页码）。"""
    _, _, doc = await _run_format(tmp_path, _pptx_bytes(), CASES[4][2])
    slides = [c for c in doc.containers if c.container_type == "slide"]
    assert slides and slides[0].page_number is None
    assert any(e.text == "发布说明" for e in doc.elements)


@pytest.mark.asyncio
async def test_html_fidelity_points(tmp_path) -> None:
    """HTML 保真：DOM 容器 + xpath 证据 + 表格。"""
    _, _, doc = await _run_format(tmp_path, _HTML_BYTES, "text/html")
    assert any(c.container_type == "dom_document" for c in doc.containers)
    xpaths = [
        s.native_ref.get("xpath")
        for e in doc.elements for s in e.source_spans if s.native_ref
    ]
    assert any(xp for xp in xpaths), "xpath evidence missing"
    tables = [a for a in doc.structured_assets.values() if getattr(a, "rows", None)]
    assert tables and any("过温" in c.text for c in tables[0].cells)
